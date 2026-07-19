import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

def hex_to_rgb(hex_code: str) -> RGBColor:
    hex_code = hex_code.lstrip('#')
    if len(hex_code) == 3:
        hex_code = ''.join(c * 2 for c in hex_code)
    return RGBColor(int(hex_code[0:2], 16), int(hex_code[2:4], 16), int(hex_code[4:6], 16))

# ─── Theme Map: all 12 themes → bg, card, text, muted, primary, secondary ────
THEME_COLORS = {
    'aurora':         {'bg': '0f172a', 'card': '1e293b', 'text': 'f0f9ff', 'muted': '94a3b8', 'primary': '22d3ee', 'secondary': '38bdf8'},
    'midnight_vc':    {'bg': '0d0d1a', 'card': '1a1a2e', 'text': 'e2e8f0', 'muted': '64748b', 'primary': 'a78bfa', 'secondary': '818cf8'},
    'deep_ocean':     {'bg': '020617', 'card': '0f172a', 'text': 'bfdbfe', 'muted': '475569', 'primary': '38bdf8', 'secondary': '0ea5e9'},
    'obsidian_gold':  {'bg': '0c0a09', 'card': '1c1917', 'text': 'fef3c7', 'muted': '78716c', 'primary': 'f59e0b', 'secondary': 'd97706'},
    'crimson_empire': {'bg': '0a0000', 'card': '1a0a0a', 'text': 'fee2e2', 'muted': '78716c', 'primary': 'ef4444', 'secondary': 'dc2626'},
    'forest_venture': {'bg': '042f1a', 'card': '064e2f', 'text': 'd1fae5', 'muted': '6b7280', 'primary': '10b981', 'secondary': '059669'},
    'silicon_sage':   {'bg': '0f1a12', 'card': '1a2e1c', 'text': 'dcfce7', 'muted': '6b7280', 'primary': '4ade80', 'secondary': '22c55e'},
    'arctic_light':   {'bg': 'f8fafc', 'card': 'ffffff', 'text': '0f172a', 'muted': '64748b', 'primary': '0ea5e9', 'secondary': '38bdf8'},
    'corporate_slate':{'bg': '0f172a', 'card': '1e293b', 'text': 'f1f5f9', 'muted': '94a3b8', 'primary': '3b82f6', 'secondary': '60a5fa'},
    'rose_quartz':    {'bg': '18000d', 'card': '2d0019', 'text': 'fce7f3', 'muted': '9ca3af', 'primary': 'f472b6', 'secondary': 'ec4899'},
    'minimal_light':  {'bg': 'ffffff', 'card': 'f8fafc', 'text': '0f172a', 'muted': '64748b', 'primary': '6366f1', 'secondary': '818cf8'},
    'corporate_blue': {'bg': '0f172a', 'card': '1e3a8a', 'text': 'ffffff', 'muted': 'bfdbfe', 'primary': '3b82f6', 'secondary': '60a5fa'},
}

def _add_solid_bg(slide, rgb: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def _add_rect(slide, x, y, w, h, fill_rgb: RGBColor, line_rgb: RGBColor = None, line_width_pt: float = 0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb and line_width_pt > 0:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def _add_text(slide, text, x, y, w, h, size_pt, bold, rgb: RGBColor, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = bold
    p.font.size = Pt(size_pt)
    p.font.color.rgb = rgb
    p.font.name = 'Inter'
    p.alignment = align
    return txBox

def _add_eyebrow(slide, tag, x, y, w, primary_rgb: RGBColor, bg_rgb: RGBColor):
    """Add a small tag pill above the title"""
    # Background pill
    pill = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(0.3))  # Rounded rect
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg_rgb
    pill.line.fill.background()
    # Text
    tf = pill.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = tag.upper()
    p.font.bold = True
    p.font.size = Pt(9)
    p.font.color.rgb = primary_rgb
    p.font.name = 'Inter'
    p.alignment = PP_ALIGN.CENTER


def generate_pptx(deck_data, theme, output_path):
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Resolve theme colors
    C = THEME_COLORS.get(theme, THEME_COLORS['aurora'])
    bg_rgb       = hex_to_rgb(C['bg'])
    card_rgb     = hex_to_rgb(C['card'])
    text_rgb     = hex_to_rgb(C['text'])
    muted_rgb    = hex_to_rgb(C['muted'])
    primary_rgb  = hex_to_rgb(C['primary'])
    secondary_rgb = hex_to_rgb(C['secondary'])

    slides = deck_data.get('slides', [])
    if not slides:
        slides = [{'type': 'title', 'companyName': 'Company', 'tagline': 'Tagline'}]

    for slide_data in slides:
        st = slide_data.get('type', 'list')
        ps = prs.slides.add_slide(blank)
        _add_solid_bg(ps, bg_rgb)

        tag   = slide_data.get('tag', '')
        title = slide_data.get('title', '')

        # ── Accent top strip (1px)
        _add_rect(ps, 0, 0, 13.333, 0.04, primary_rgb)

        # ── Helper: add standard eyebrow + title block
        def add_header(tx, ty, tw=11.33, eyebrow_offset=0.5):
            if tag:
                _add_text(ps, tag, tx, ty, tw, 0.3, 8, True, primary_rgb, PP_ALIGN.LEFT)
                ty += eyebrow_offset
            if title:
                _add_text(ps, title, tx, ty, tw, 0.9, 32, True, text_rgb, PP_ALIGN.LEFT)
                ty += 1.05
            return ty

        # ════════════════════════════════════════
        # TITLE SLIDE
        # ════════════════════════════════════════
        if st == 'title':
            company = slide_data.get('companyName', 'Company')
            tagline = slide_data.get('tagline', '')

            # Right-side geometric accent panel
            panel = _add_rect(ps, 8.0, 0.04, 5.333, 7.46, card_rgb)

            # Grid dots simulation (6x4 dots)
            for row in range(4):
                for col in range(6):
                    dot = ps.shapes.add_shape(9, Inches(8.4 + col * 0.7), Inches(1.0 + row * 1.3), Inches(0.08), Inches(0.08))
                    dot.fill.solid()
                    dot.fill.fore_color.rgb = primary_rgb
                    dot.line.fill.background()

            # Big initial letter
            initial = company[0].upper() if company else 'C'
            _add_text(ps, initial, 9.2, 1.5, 3.5, 4.5, 110, True, muted_rgb, PP_ALIGN.CENTER)

            # Left content
            if tag:
                _add_text(ps, f'— {tag}', 1.0, 1.8, 6.5, 0.4, 9, True, primary_rgb, PP_ALIGN.LEFT)
            _add_text(ps, company, 1.0, 2.3, 6.5, 2.0, 64, True, text_rgb, PP_ALIGN.LEFT)

            # Divider line
            _add_rect(ps, 1.0, 4.5, 1.5, 0.06, primary_rgb)

            _add_text(ps, tagline, 1.0, 4.75, 6.5, 1.5, 22, False, muted_rgb, PP_ALIGN.LEFT)

        # ════════════════════════════════════════
        # SPLIT SLIDE (Problem / Solution / How)
        # ════════════════════════════════════════
        elif st == 'split':
            is_problem = 'problem' in tag.lower()
            bullet_color = hex_to_rgb('ef4444') if is_problem else primary_rgb
            bullet_char  = '✗' if is_problem else '✓'

            content = slide_data.get('content', '')
            items   = slide_data.get('items', [])

            ty = add_header(0.8, 0.5, tw=6.5)

            # Content sentence
            if content:
                _add_text(ps, content, 0.8, ty, 6.2, 1.0, 15, False, muted_rgb, PP_ALIGN.LEFT)
                ty += 1.1

            # Bullet items (left column)
            for item in items[:4]:
                _add_text(ps, f'{bullet_char}  {item}', 0.8, ty, 6.2, 0.65, 14, False, text_rgb, PP_ALIGN.LEFT)
                ty += 0.75

            # Right side image placeholder card
            img_rect = _add_rect(ps, 7.6, 0.7, 5.0, 6.2, card_rgb, primary_rgb, 1.0)
            _add_text(ps, '[ Image ]', 7.6, 3.3, 5.0, 0.6, 14, False, muted_rgb, PP_ALIGN.CENTER)

        # ════════════════════════════════════════
        # LIST SLIDE (How It Works / Team / Ask)
        # ════════════════════════════════════════
        elif st == 'list':
            is_team = 'team' in tag.lower()
            is_ask  = 'ask' in tag.lower()
            items   = slide_data.get('items', [])
            content = slide_data.get('content', '')

            if is_team:
                # 2x2 Team Card Grid
                _add_text(ps, tag, 0.8, 0.3, 11.33, 0.3, 8, True, primary_rgb, PP_ALIGN.LEFT)
                _add_text(ps, title, 0.8, 0.65, 11.33, 0.9, 30, True, text_rgb, PP_ALIGN.LEFT)

                positions = [(0.5, 2.0), (7.0, 2.0), (0.5, 4.6), (7.0, 4.6)]
                accent_hex = [C['primary'], 'f59e0b', '10b981', '6366f1']

                for i, item in enumerate(items[:4]):
                    if i >= 4: break
                    cx, cy = positions[i]
                    acol = hex_to_rgb(accent_hex[i % 4])

                    # Card background
                    _add_rect(ps, cx, cy, 6.0, 2.2, card_rgb, acol, 0.5)

                    # Parse "Name · Title — Credential"
                    dot_parts = item.split('·')
                    name_part = dot_parts[0].replace('Advisor:', '').strip()
                    rest = dot_parts[1] if len(dot_parts) > 1 else item
                    dash_parts = rest.split('—')
                    title_part = dash_parts[0].strip()
                    cred_part  = dash_parts[1].strip() if len(dash_parts) > 1 else ''

                    # Avatar circle (simulated with colored rect)
                    initials = ''.join(w[0] for w in name_part.split()[:2]).upper()
                    av = ps.shapes.add_shape(9, Inches(cx + 0.15), Inches(cy + 0.5), Inches(0.9), Inches(0.9))
                    av.fill.solid(); av.fill.fore_color.rgb = acol; av.line.fill.background()
                    atf = av.text_frame; atf.paragraphs[0].text = initials
                    atf.paragraphs[0].font.bold = True; atf.paragraphs[0].font.size = Pt(16)
                    atf.paragraphs[0].font.color.rgb = hex_to_rgb(C['bg'])
                    atf.paragraphs[0].alignment = PP_ALIGN.CENTER

                    _add_text(ps, name_part,   cx + 1.2, cy + 0.35, 4.6, 0.35, 14, True,  text_rgb, PP_ALIGN.LEFT)
                    _add_text(ps, title_part,  cx + 1.2, cy + 0.72, 4.6, 0.3,  11, True,  acol,     PP_ALIGN.LEFT)
                    _add_text(ps, cred_part,   cx + 0.15, cy + 1.5,  5.7, 0.55, 10, False, muted_rgb,PP_ALIGN.LEFT)

            elif is_ask:
                # Extract dollar amount from title
                import re
                dollar_match = re.search(r'[\$£€][\d.,]+[KMBkm]?', title)
                amount = dollar_match.group(0) if dollar_match else '$—'

                # Left: items list
                ty = add_header(0.8, 0.5, tw=7.5)
                if content:
                    _add_text(ps, content, 0.8, ty, 7.2, 0.5, 14, False, muted_rgb, PP_ALIGN.LEFT)
                    ty += 0.7
                for item in items:
                    pct_match = re.search(r'(\d+)%', item)
                    _add_text(ps, f'→  {item}', 0.8, ty, 7.2, 0.6, 13, False, text_rgb, PP_ALIGN.LEFT)
                    ty += 0.7

                # Right: big amount panel
                panel = _add_rect(ps, 9.0, 0.7, 3.8, 6.1, card_rgb, primary_rgb, 1.0)
                _add_text(ps, amount,    9.0, 2.2, 3.8, 1.5, 48, True, primary_rgb, PP_ALIGN.CENTER)
                _add_text(ps, 'SEEKING', 9.0, 3.7, 3.8, 0.4, 10, True, muted_rgb,   PP_ALIGN.CENTER)

                # Allocation bars
                bar_y = 4.3
                for item in items:
                    pct_match = re.search(r'(\d+)%', item)
                    if pct_match:
                        pct = int(pct_match.group(1))
                        label = item.split('%')[0].strip().split(' ')[-1] if '%' in item else ''
                        _add_rect(ps, 9.2, bar_y, 3.4 * pct / 100, 0.18, primary_rgb)
                        _add_rect(ps, 9.2 + 3.4 * pct / 100, bar_y, 3.4 * (1 - pct/100), 0.18, hex_to_rgb(C['card']))
                        _add_text(ps, f'{pct}%', 9.2, bar_y + 0.18, 3.4, 0.3, 9, False, muted_rgb, PP_ALIGN.LEFT)
                        bar_y += 0.55

            else:
                # Standard HOW IT WORKS list
                ty = add_header(0.8, 0.5)
                if content:
                    _add_text(ps, content, 0.8, ty, 11.33, 0.6, 15, False, muted_rgb, PP_ALIGN.LEFT)
                    ty += 0.75
                for i, item in enumerate(items[:5]):
                    card = _add_rect(ps, 0.8, ty, 11.73, 0.8, card_rgb, primary_rgb, 0)
                    # left border accent
                    _add_rect(ps, 0.8, ty, 0.07, 0.8, primary_rgb)
                    _add_text(ps, item, 1.1, ty + 0.1, 11.4, 0.6, 14, False, text_rgb, PP_ALIGN.LEFT)
                    ty += 1.0

        # ════════════════════════════════════════
        # METRICS SLIDE (Traction / Business Model)
        # ════════════════════════════════════════
        elif st == 'metrics':
            metrics = slide_data.get('metrics', [])
            content = slide_data.get('content', '')

            _add_text(ps, tag,   0.8, 0.35, 11.73, 0.3,  8,  True, primary_rgb, PP_ALIGN.CENTER)
            _add_text(ps, title, 0.8, 0.65, 11.73, 0.85, 30, True, text_rgb,    PP_ALIGN.CENTER)
            if content:
                _add_text(ps, content, 0.8, 1.5, 11.73, 0.5, 14, False, muted_rgb, PP_ALIGN.CENTER)

            n = len(metrics)
            if n == 0:
                continue

            card_w = min(11.0 / n - 0.3, 3.2)
            total_w = n * card_w + (n - 1) * 0.4
            start_x = (13.333 - total_w) / 2

            for i, m in enumerate(metrics):
                cx = start_x + i * (card_w + 0.4)
                # Card with top accent border
                _add_rect(ps, cx, 2.2, card_w, 0.1, primary_rgb)        # top accent
                _add_rect(ps, cx, 2.3, card_w, 3.5, card_rgb, None, 0)  # card body

                # Value (big number)
                _add_text(ps, m.get('value', '-'), cx, 2.5, card_w, 1.2,
                          38, True, primary_rgb, PP_ALIGN.CENTER)
                # Label
                _add_text(ps, m.get('label', ''), cx, 3.7, card_w, 0.5,
                          11, True, muted_rgb, PP_ALIGN.CENTER)
                # Desc
                if m.get('desc'):
                    _add_text(ps, m['desc'], cx, 4.2, card_w, 0.5,
                              10, False, muted_rgb, PP_ALIGN.CENTER)

        # ════════════════════════════════════════
        # CHART SLIDE — data table fallback
        # ════════════════════════════════════════
        elif st == 'chart':
            chart_type = slide_data.get('chartType', 'bar')
            data       = slide_data.get('data', [])
            content    = slide_data.get('content', '')

            ty = add_header(0.8, 0.5)
            if content:
                _add_text(ps, content, 0.8, ty, 11.73, 0.5, 14, False, muted_rgb, PP_ALIGN.LEFT)
                ty += 0.65

            if not data:
                _add_text(ps, 'No data', 0.8, ty, 11.73, 1.0, 20, False, muted_rgb, PP_ALIGN.CENTER)
            else:
                # Render as visual bars
                max_val = max(d.get('value', 1) for d in data) or 1
                bar_area_w = 10.5
                bar_h_each = min(0.95, (7.5 - ty - 0.5) / len(data) - 0.15)

                for i, d in enumerate(data):
                    val   = d.get('value', 0)
                    label = d.get('label', '')
                    pct   = val / max_val
                    bar_w = pct * (bar_area_w - 2.5)

                    # Label
                    _add_text(ps, label, 0.8, ty + 0.05, 2.3, bar_h_each, 11, False, muted_rgb, PP_ALIGN.RIGHT)

                    # Background track
                    _add_rect(ps, 3.3, ty + 0.1, bar_area_w - 2.5, bar_h_each - 0.1, card_rgb)

                    # Value bar (color cycle)
                    accent_cycle = [C['primary'], 'f59e0b', '10b981', '6366f1', 'f43f5e', '84cc16']
                    bar_rgb = hex_to_rgb(accent_cycle[i % len(accent_cycle)])
                    if bar_w > 0.01:
                        _add_rect(ps, 3.3, ty + 0.1, bar_w, bar_h_each - 0.1, bar_rgb)

                    # Value label at end
                    val_str = f'{val}' if isinstance(val, int) else f'{val:.1f}'
                    _add_text(ps, val_str, 3.3 + bar_w + 0.05, ty + 0.05, 1.5, bar_h_each, 11, True, primary_rgb, PP_ALIGN.LEFT)

                    ty += bar_h_each + 0.18

    prs.save(output_path)
    return True
